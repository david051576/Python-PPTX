import os
import re
import json
import urllib2
import requests
import cookielib
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR



#  This program will detect the types your page datas and automatically generate the corresponding pattern of them.


"""
Just follow the pattern below, the less the number is, the older the datas should be.

Pattern #1 Important quote
               (1) Quote you want to emphasize
Pattern #2 Issue with photo
               (1) Title you want to include 
               (2) Paragraphs you want to include
               (3) Photo's keywords (you can put the title's name)
Pattern #3 Two concepts
               (1) Title you want to include 
               (2) Concept one
               (3) Paragragh one
               (4) Concept two
               (5) Paragraph two
               
# Interpret JSON in this program??
"""


def get_soup(url,header):  # get pictures from the web

    return BeautifulSoup(urllib2.urlopen(urllib2.Request(url,headers=header)),'html.parser')


def quoto(pagedata): # add quoto in the middle of the slide
    
    # Slide information    
    frame_width = Inches(10)
    frame_height = Inches(7.5)


    # Quote Paragraph added
    left = frame_width/5.0
    top = frame_height/4.0
    width = frame_width/5.0*3.0
    height = frame_width/5.0*2.0

    # Textbox Position
    textBox = slide.shapes.add_textbox(left, top, width, height)
    tf1 = textBox.text_frame
    tf1.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Text Format
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = "'" + pagedata + "'"
    font1 = run1.font
    font1.name = 'Calibri'
    font1.size = Pt(48)
    font1.bold = True
    font1.italic = True
    font1.color.rgb = RGBColor(0x00, 0x00, 0x00)


pic_num = 1 # how many pictures do you want to download


##### Download the related images

##### Step1  pre-process before google searching

query = raw_input("input the quoto: ") # the keywords for google to search for
sentence = query
query= query.split()  # pre-process
query='+'.join(query)
url="https://www.google.co.in/search?q="+query+"&source=lnms&tbm=isch"  # google format of searching
print url


##### Step2  Create the file, refresh the dictionary

DIR="Pictures"  # add the directory name for your image
header={'User-Agent':"Mozilla/5.0 (Windows NT 6.1; WOW64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/43.0.2357.134 Safari/537.36"}
soup = get_soup(url,header)


##### Step3  Calcutlate total images on the searching page

ActualImages=[] # contains the link for Large original images, type of  image
for a in soup.find_all("div",{"class":"rg_meta"}):
    link , Type =json.loads(a.text)["ou"]  ,json.loads(a.text)["ity"]
    ActualImages.append((link,Type))
    if len(ActualImages) > pic_num-1:
        break

print  "total" , len(ActualImages),"images will be downloaded."

if not os.path.exists(DIR):
            os.mkdir(DIR)
DIR = os.path.join(DIR, query.split()[0])

if not os.path.exists(DIR):
            os.mkdir(DIR)


##### Step4  Download the images
            
for i , (img , Type) in enumerate( ActualImages):
    
    try:
        req = urllib2.Request(img, headers={'User-Agent' : header})
        raw_img = urllib2.urlopen(req).read()

        cntr = len([i for i in os.listdir(DIR) if query in i]) + 1
        print cntr

        if len(Type)==0:
            f = open(os.path.join(DIR , query+  "_"+ str(cntr)+".jpg"), 'wb')
        else :
            f = open(os.path.join(DIR , query + "_"+ str(cntr)+".jpg"), 'wb')


        f.write(raw_img)
        f.close()

    except Exception as e:
        print "could not load : "+img
        print e



##### Create new PowerPoint file and select the slide pattern
        
prs = Presentation()

slide = prs.slides.add_slide(prs.slide_layouts[0])
shapes = slide.shapes
shapes.title.text = "GOOD PRESENTATION in Microsoft!!"

title_only_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes
            
top = 0
left = 0
img_path = os.path.join(DIR , query + "_"+ str(cntr)+".jpg")
pic = shapes.add_picture(img_path, left, top, Inches(10), Inches(7.5))


shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(10), Inches(7.5))
            
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xC0, 0xC0, 0xC5)
fill.fore_color.brightness = 0.25
fill.transparency = 0.75
    
quoto(sentence) # Quoto program

# Save the file name
prs.save('microsoft.pptx')





